"""
存储服务
"""
import subprocess

import io
import os
import shutil
from typing import   Optional, List, Any
from docx import Document
from pptx import Presentation

from minio import Minio
temp_dir = "/tmp/libreoffice_convert"
os.makedirs(temp_dir, exist_ok=True)
EXTS = ["doc", "docx", "ppt", "pptx"]

OSS_APIKEY="FnW30Im6qnfbXMHjRBDV"
OSS_APIKEY_SECRETE="r2HX3b5zfLAvaKgFYk9F9Zj9sFbeTbqO1PWBNljz"
OSS_END_POINT="192.168.81.210:9010"
OSS_BUCKET_NAME="corpus-origin"

import datetime as dt

from pydantic import BaseModel, validator
# OSS_APIKEY="kbNFYDw3CnUL5Wzj"
# OSS_APIKEY_SECRETE="cVstAkFEIMjFmfIZZ8uztqZ1CYzhvYQZ"
# OSS_END_POINT="139.226.106.165:9000"
# OSS_BUCKET_NAME="kps"

class Model(BaseModel):
    @validator("*")
    def add_tzinfo(cls, v):
        if isinstance(v, dt.datetime) and not v.tzinfo:
            return v.replace(tzinfo=dt.timezone.utc)
        return v

    def json(
        self,
        ensure_ascii: bool = False,
        **kwargs,
    ) -> str:
        return super().json(ensure_ascii=ensure_ascii, **kwargs)

class FileDirectoryOutput(Model):
    file_name: str
    file_size: Optional[int]
    is_directory: bool = True

class S3Service:
    def __init__(self) -> None:
        self._endpoint_url = OSS_END_POINT
        self._oss_bucket_name = OSS_BUCKET_NAME
        self._aws_access_key_id = OSS_APIKEY
        self._aws_secret_access_key = OSS_APIKEY_SECRETE
        self.client = Minio(self._endpoint_url, access_key=self._aws_access_key_id,
                            secret_key=self._aws_secret_access_key, secure=False)

    def upload_file_to_minio(self, file_path: str, key: str):
        print("开始上传文件到Minio...")

        self.client.fput_object(self._oss_bucket_name, key, file_path)
        print("开始上传文件到Minio结束...")
    def upload_file_by_minio(self, file_path: str, key: str):
        with open(file_path, 'rb') as f:
            file_bytes = f.read()
            self.upload_file_bytes_by_minio(file_bytes, key=key)

    def upload_file_bytes_by_minio(self, file_bytes: bytes, key: str) -> dict:
        """上传 对象，返回 url 等信息"""

        self.client.put_object(self._oss_bucket_name, key, io.BytesIO(file_bytes), len(file_bytes))

        return {"msg": f'{key}上传成功',
                "store_key": f"{key}"}

    def get_file_by_minio(self, key: str,temp_file:str) ->str:
        respone = self.client.fget_object(self._oss_bucket_name, key, temp_file)
        return respone.object_name
    def list_files(self, prefix: Optional[str] = None) -> List[FileDirectoryOutput]:
        result = []
        if prefix is None:
            return result
        try:
            response = self.client.list_objects(self._oss_bucket_name)
            print(response)
            for i in response['Contents']:
                print(i)
                print("Name: ", i['Key'])
        except Exception as e:
            return result

    def list_dirs(self, prefix: Optional[str] = None) -> List[FileDirectoryOutput]:
        dirs = []
        objects = self.client.list_objects(self._oss_bucket_name, prefix=prefix, recursive=False)

        for obj in objects:
            if obj.is_dir:
                dirs.append(FileDirectoryOutput(file_name=obj.object_name, file_size=obj.size, is_directory=True))
            else:
                dirs.append(FileDirectoryOutput(file_name=obj.object_name, file_size=obj.size, is_directory=False))
        return dirs
 # 临时工作目录

def clean_fonts(input,output,file_type):
    if file_type.lower() in ["docx"]:
        wordcleantfont(input,output)
    if file_type.lower() in ["doc"]:
        def convert_doc_to_docx(input, output, file_type):
            outdir, _ = os.path.split(input)
            try:
                subprocess.run([
                    "libreoffice",
                    "--headless",            # 无界面
                    "--convert-to", "docx",   # 指定转换格式
                    "--outdir", outdir,   # 指定输出目录
                    input
                ], check=True)
                print(f"成功转换: {input}")
                if file_type == "doc":
                    wordcleantfont(input+'x', output+'x') 
                else:
                    wordcleantfont(input+'X', output+'X')
                
            except subprocess.CalledProcessError as e:
                print(f"转换失败: {input}")
                print(e)
                shutil.copyfile(input, output)
                # wordcleantfont(input, output)
                
        convert_doc_to_docx(input, output, file_type)
        
    if file_type.lower() in ["ppt","pptx"]:
        ppt_clean_fonts(input,output)
        
def wordcleantfont(input,output):

    doc = Document(input)

    # 默认字体（例如：'Calibri'）
    default_font = 'WenQuanYi Zen Hei'

    # 遍历文档中的每个段落
    for paragraph in doc.paragraphs:
        # 遍历段落中的每个run（可以包含不同的样式和字体）
        for run in paragraph.runs:
            # 设置run的字体为默认字体
            run.font.name = default_font
    doc.save(output)

def ppt_clean_fonts(input,output):

    def change_font_in_presentation(presentation_path, default_font_name='WenQuanYi Zen Hei'):
        prs = Presentation(presentation_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # 确保shape有text属性
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = default_font_name
                            # 也可以设置字体大小，例如：run.font.size = Pt(24)
                            # 设置字体颜色，例如：run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        prs.save(output)

    # 使用示例
    change_font_in_presentation(input)

class Process:
    def __init__(self) -> None:
        self.service = S3Service()
        self.file_types = {}
        self.nums = 0
        self.file_size = 0

    def run(self,files:List[Any],prefix:str):
        if len(files)==0:
            return
        print(files)
        for _index,file in enumerate(files):
            if file.file_name != prefix or not file.is_directory:
                if  not file.is_directory:
                    file_type=file.file_name.split(".")[-1]
                    print(file_type)
                    # file_name = file.file_name
                    if file_type not in self.file_types.keys():
                        self.file_types[file_type]=0

                    self.file_types[file_type]+=1
                    self.nums=self.nums+1
                    self.file_size=self.file_size+file.file_size

                    if file_type.lower() in EXTS:
                        print(file.file_name)
                        try:
                            temp_file = os.path.join(temp_dir, "temp." + file_type)
                            print("temp_file",temp_file)
                            # s3 下载
                            self.service.get_file_by_minio(file.file_name,temp_file)
                            print("下载完成")

                            print("开始转换word字体")
                            temp1_file = os.path.join(temp_dir, "1temp." + file_type)
                            clean_fonts(temp_file,temp1_file,file_type)
                            print("完成转换word字体")
                            subprocess.run([
                                "libreoffice",
                                "--headless",
                                "--convert-to", "pdf",
                                temp1_file,
                                "--outdir", temp_dir
                            ], check=True)
                            print("转换完成")
                            temp_pdf = os.path.join(temp_dir, "1temp.pdf")
                            d, f = os.path.split(file.file_name)
                            f = "kps_converted_" + f + ".pdf"
                            key_ = f"{os.path.join(d, f)}"
                            # key_=f"{file.file_name.replace(file_type, 'pdf')}"
                            print("开始上传",key_)
                            print(temp_pdf)
                            self.service.upload_file_to_minio(temp_pdf, key=key_)

                        except Exception as e:
                            print(f"报错: {e}")
                else:
                    self.run(self.service.list_dirs(file.file_name),file.file_name)

if __name__ == '__main__':
    files = S3Service().list_dirs()

    for i in files:
        if  i.is_directory is False:
            continue
        p = Process()
        p.run(files=S3Service().list_dirs(i.file_name),prefix=i.file_name)
    
    print("Done!")